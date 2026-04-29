"""Document processing — extraction, chunking, embedding orchestration.

Pipeline:
  upload → extract_text (dispatcher by mime) → chunk → embed → DB

Extractors return ExtractionResult with per-page text, image_refs, and bytes
of any embedded images. Images get stored as DocumentImage rows so the model
can later request base64 via the get_document_image tool.

Archives are NOT extracted here — they're marked sandbox_only and the model
unpacks them itself in /workspace/files/.
"""
from __future__ import annotations

import asyncio
import hashlib
import logging
import re
import uuid as _uuid
from dataclasses import dataclass, field
from pathlib import Path
from uuid import UUID

from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from quip.models.file import File, DocumentChunk, DocumentImage
from quip.services.config import get_setting

logger = logging.getLogger(__name__)


# ---------- Result types ----------

@dataclass
class PageContent:
    text: str
    page: int | None = None        # 1-based page/slide/sheet index
    image_refs: list[str] = field(default_factory=list)
    source: str = "text"           # "text" | "ocr"


@dataclass
class ExtractedImage:
    ref: str                       # marker like "img_1"
    page: int | None
    mime: str
    data: bytes


@dataclass
class ExtractionResult:
    pages: list[PageContent] = field(default_factory=list)
    images: list[ExtractedImage] = field(default_factory=list)
    sandbox_only: bool = False     # archives — no chunks except a marker


# ---------- Mime classification ----------

ARCHIVE_MIMES = {
    "application/zip",
    "application/x-zip-compressed",
    "application/x-tar",
    "application/gzip",
    "application/x-gzip",
    "application/x-rar-compressed",
    "application/vnd.rar",
    "application/x-7z-compressed",
}


def is_archive(content_type: str) -> bool:
    return content_type in ARCHIVE_MIMES


# ---------- Public entry ----------

async def extract(file_path: str, content_type: str) -> ExtractionResult:
    """Dispatcher — pick extractor by mime. Returns ExtractionResult."""
    path = Path(file_path)
    ct = (content_type or "").lower()

    if is_archive(ct):
        return ExtractionResult(sandbox_only=True)

    if ct in ("text/plain", "text/markdown", "text/csv") or ct.startswith("text/"):
        try:
            return ExtractionResult(pages=[PageContent(text=path.read_text(encoding="utf-8", errors="replace"))])
        except Exception as e:
            logger.error(f"text read failed: {e}")
            return ExtractionResult()

    if ct == "application/pdf":
        return await _extract_pdf(path)

    if ct == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return await asyncio.to_thread(_extract_docx, path)

    if ct in (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
    ):
        return await asyncio.to_thread(_extract_xlsx, path)

    if ct == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return await asyncio.to_thread(_extract_pptx, path)

    if ct == "application/epub+zip":
        return await asyncio.to_thread(_extract_epub, path)

    if ct.startswith("image/"):
        return await _extract_image(path, ct)

    # Fallback — try as text
    try:
        return ExtractionResult(pages=[PageContent(text=path.read_text(encoding="utf-8", errors="replace"))])
    except Exception:
        return ExtractionResult()


# ---------- PDF ----------

async def _extract_pdf(path: Path) -> ExtractionResult:
    """PyMuPDF text + embedded images, with cloud OCR fallback for scan-like pages."""
    try:
        import fitz  # pymupdf
    except ImportError:
        logger.warning("pymupdf not installed")
        return ExtractionResult()

    pages: list[PageContent] = []
    images: list[ExtractedImage] = []
    img_counter = 0
    ocr_pages: list[tuple[int, bytes]] = []  # (page_index, png bytes)

    try:
        doc = fitz.open(str(path))
    except Exception as e:
        logger.error(f"PDF open failed: {e}")
        return ExtractionResult()

    try:
        for pi, page in enumerate(doc):
            text = page.get_text() or ""

            # Embedded images
            page_refs: list[str] = []
            for img in page.get_images(full=True):
                xref = img[0]
                try:
                    base = doc.extract_image(xref)
                    img_counter += 1
                    ref = f"img_{img_counter}"
                    images.append(ExtractedImage(
                        ref=ref,
                        page=pi + 1,
                        mime=f"image/{base['ext']}",
                        data=base["image"],
                    ))
                    page_refs.append(ref)
                except Exception:
                    continue

            # OCR trigger: very little text OR images cover >60% of page
            needs_ocr = False
            if len(text.strip()) < 20:
                needs_ocr = True
            else:
                try:
                    page_area = page.rect.width * page.rect.height
                    if page_area > 0:
                        img_area = sum(
                            (info["bbox"][2] - info["bbox"][0]) * (info["bbox"][3] - info["bbox"][1])
                            for info in page.get_image_info()
                        )
                        if img_area / page_area > 0.6:
                            needs_ocr = True
                except Exception:
                    pass

            if needs_ocr:
                try:
                    pix = page.get_pixmap(dpi=150)
                    ocr_pages.append((pi, pix.tobytes("png")))
                except Exception:
                    pass

            marker_suffix = (" " + " ".join(f"[image: {r}]" for r in page_refs)) if page_refs else ""
            pages.append(PageContent(
                text=text + marker_suffix,
                page=pi + 1,
                image_refs=page_refs,
                source="text",
            ))
    finally:
        doc.close()

    # OCR fallback — replace text on pages that look scanned
    if ocr_pages:
        from quip.services.ocr import ocr_image
        for pi, png_bytes in ocr_pages:
            res = await ocr_image(png_bytes, "image/png")
            if res.text:
                pages[pi].text = res.text + (
                    " " + " ".join(f"[image: {r}]" for r in pages[pi].image_refs)
                    if pages[pi].image_refs else ""
                )
                pages[pi].source = "ocr"

    return ExtractionResult(pages=pages, images=images)


# ---------- DOCX ----------

def _extract_docx(path: Path) -> ExtractionResult:
    try:
        from docx import Document
    except ImportError:
        logger.warning("python-docx not installed")
        return ExtractionResult()

    try:
        doc = Document(str(path))
    except Exception as e:
        logger.error(f"DOCX open failed: {e}")
        return ExtractionResult()

    parts: list[str] = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)

    for tbl in doc.tables:
        rows = []
        for row in tbl.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
        if rows:
            parts.append("\n".join(rows))

    images: list[ExtractedImage] = []
    img_counter = 0
    refs: list[str] = []
    try:
        for rel in doc.part.rels.values():
            if "image" in (rel.reltype or ""):
                try:
                    blob = rel.target_part.blob
                    ct = rel.target_part.content_type or "image/png"
                    img_counter += 1
                    ref = f"img_{img_counter}"
                    images.append(ExtractedImage(ref=ref, page=None, mime=ct, data=blob))
                    refs.append(ref)
                except Exception:
                    continue
    except Exception:
        pass

    text = "\n\n".join(parts)
    if refs:
        text += "\n\n" + " ".join(f"[image: {r}]" for r in refs)
    return ExtractionResult(
        pages=[PageContent(text=text, image_refs=refs)],
        images=images,
    )


# ---------- XLSX ----------

XLSX_MAX_ROWS_INLINE = 200
XLSX_MAX_COLS = 50
XLSX_SAMPLE_HEAD = 50
XLSX_SAMPLE_TAIL = 50


def _extract_xlsx(path: Path) -> ExtractionResult:
    try:
        from openpyxl import load_workbook
    except ImportError:
        logger.warning("openpyxl not installed")
        return ExtractionResult()

    try:
        wb = load_workbook(str(path), read_only=True, data_only=True)
    except Exception as e:
        logger.error(f"XLSX open failed: {e}")
        return ExtractionResult()

    pages: list[PageContent] = []
    for sheet_idx, ws in enumerate(wb.worksheets):
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        # truncate columns
        rows = [r[:XLSX_MAX_COLS] for r in rows]
        n = len(rows)

        if n <= XLSX_MAX_ROWS_INLINE:
            text = _rows_to_md(rows)
        else:
            head = rows[:XLSX_SAMPLE_HEAD]
            tail = rows[-XLSX_SAMPLE_TAIL:]
            text = (
                _rows_to_md(head)
                + f"\n\n[truncated: {n - XLSX_SAMPLE_HEAD - XLSX_SAMPLE_TAIL} more rows]\n\n"
                + _rows_to_md(tail)
                + "\n\n"
                + _column_summary(rows)
            )
        text = f"## Sheet: {ws.title}\n\n" + text
        pages.append(PageContent(text=text, page=sheet_idx + 1))

    wb.close()
    return ExtractionResult(pages=pages)


def _rows_to_md(rows: list[tuple]) -> str:
    out = []
    for r in rows:
        cells = ["" if c is None else str(c).replace("\n", " ").replace("|", "\\|") for c in r]
        out.append("| " + " | ".join(cells) + " |")
    return "\n".join(out)


def _column_summary(rows: list[tuple]) -> str:
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    summaries = []
    for col in range(width):
        nums = []
        for r in rows:
            if col < len(r) and isinstance(r[col], (int, float)):
                nums.append(r[col])
        if nums and len(nums) >= 5:
            summaries.append(
                f"col {col}: n={len(nums)}, min={min(nums)}, max={max(nums)}, mean={sum(nums)/len(nums):.2f}"
            )
    return "[column stats]\n" + "\n".join(summaries) if summaries else ""


# ---------- PPTX ----------

def _extract_pptx(path: Path) -> ExtractionResult:
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed")
        return ExtractionResult()

    try:
        prs = Presentation(str(path))
    except Exception as e:
        logger.error(f"PPTX open failed: {e}")
        return ExtractionResult()

    pages: list[PageContent] = []
    images: list[ExtractedImage] = []
    img_counter = 0

    for si, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        slide_refs: list[str] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in p.runs)
                    if line.strip():
                        parts.append(line)
            if getattr(shape, "shape_type", None) == 13:  # PICTURE
                try:
                    img = shape.image
                    img_counter += 1
                    ref = f"img_{img_counter}"
                    images.append(ExtractedImage(
                        ref=ref, page=si, mime=img.content_type, data=img.blob,
                    ))
                    slide_refs.append(ref)
                except Exception:
                    continue

        # Notes
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    parts.append("[notes]\n" + notes)
        except Exception:
            pass

        text = f"## Slide {si}\n\n" + "\n".join(parts)
        if slide_refs:
            text += "\n\n" + " ".join(f"[image: {r}]" for r in slide_refs)
        pages.append(PageContent(text=text, page=si, image_refs=slide_refs))

    return ExtractionResult(pages=pages, images=images)


# ---------- EPUB ----------

def _extract_epub(path: Path) -> ExtractionResult:
    try:
        from ebooklib import epub, ITEM_DOCUMENT
        from bs4 import BeautifulSoup
    except ImportError:
        logger.warning("ebooklib/beautifulsoup4 not installed")
        return ExtractionResult()

    try:
        book = epub.read_epub(str(path))
    except Exception as e:
        logger.error(f"EPUB open failed: {e}")
        return ExtractionResult()

    pages: list[PageContent] = []
    chapter_idx = 0
    for item in book.get_items_of_type(ITEM_DOCUMENT):
        try:
            html = item.get_content().decode("utf-8", errors="replace")
            soup = BeautifulSoup(html, "html.parser")
            text = soup.get_text(separator="\n").strip()
        except Exception:
            continue
        if not text:
            continue
        chapter_idx += 1
        pages.append(PageContent(text=text, page=chapter_idx))

    return ExtractionResult(pages=pages)


# ---------- Standalone image ----------

async def _extract_image(path: Path, mime: str) -> ExtractionResult:
    """OCR a standalone uploaded image into RAG chunks."""
    try:
        data = path.read_bytes()
    except Exception as e:
        logger.error(f"image read failed: {e}")
        return ExtractionResult()

    from quip.services.ocr import ocr_image
    res = await ocr_image(data, mime)
    if not res.text or not res.text.strip():
        return ExtractionResult()

    return ExtractionResult(pages=[PageContent(text=res.text, source="ocr")])


# ---------- Chunking (preserves page provenance) ----------

def chunk_pages(
    pages: list[PageContent],
    max_tokens: int = 512,
    overlap_tokens: int = 64,
) -> list[tuple[str, dict, str]]:
    """Chunk pages into (text, metadata, content_hash) tuples.

    Metadata carries page, source, image_refs. content_hash is SHA-256 hex
    of the chunk text — used for cross-file dedup in RAG retrieval.
    """
    if not pages:
        return []

    try:
        import tiktoken
        enc = tiktoken.get_encoding("cl100k_base")
        count_tokens = lambda t: len(enc.encode(t))
    except ImportError:
        count_tokens = lambda t: len(t) // 4

    out: list[tuple[str, dict, str]] = []
    for page in pages:
        if not page.text.strip():
            continue
        chunks = _chunk_one(page.text, max_tokens, overlap_tokens, count_tokens)
        for chunk in chunks:
            refs_in_chunk = [r for r in page.image_refs if f"[image: {r}]" in chunk]
            meta = {"source": page.source}
            if page.page is not None:
                meta["page"] = page.page
            if refs_in_chunk:
                meta["image_refs"] = refs_in_chunk
            chash = hashlib.sha256(chunk.encode()).hexdigest()
            out.append((chunk, meta, chash))
    return out


_HEADING_RE = re.compile(r'^#{1,4}\s+.+$', re.MULTILINE)
_SECTION_BREAK_RE = re.compile(r'\n{2,}')


def _chunk_one(text: str, max_tokens: int, overlap_tokens: int, count_tokens) -> list[str]:
    """Split text into chunks respecting structural boundaries.

    Strategy:
    1. Split on double-newlines (paragraph/section breaks) first.
    2. Within oversized paragraphs, fall back to sentence-level splitting.
    3. Markdown headings anchor new chunks so a heading isn't orphaned from
       its body text.
    """
    # Split on section boundaries first
    sections = _SECTION_BREAK_RE.split(text)
    # Further split long sections on headings
    segments: list[str] = []
    for sec in sections:
        sec = sec.strip()
        if not sec:
            continue
        # Split on heading boundaries but keep heading with following text
        heading_splits = _HEADING_RE.split(sec)
        if len(heading_splits) > 1:
            # Re-attach headings — re.finditer gives us the actual heading lines
            headings: list[str] = [m.group(0) for m in _HEADING_RE.finditer(sec)]
            for i, part in enumerate(heading_splits):
                part = part.strip()
                if not part:
                    continue
                if i < len(headings):
                    segments.append(headings[i] + "\n" + part)
                else:
                    segments.append(part)
        else:
            segments.append(sec)

    # Sentence-level splitting for segments that are still too large
    fine_segments: list[str] = []
    for seg in segments:
        tok = count_tokens(seg)
        if tok <= max_tokens:
            fine_segments.append(seg)
        else:
            parts = re.split(r'(\.\s|!\s|\?\s|;\s|:\s)', seg)
            for i in range(0, len(parts), 2):
                s = parts[i]
                if i + 1 < len(parts):
                    s += parts[i + 1]
                if s.strip():
                    fine_segments.append(s)

    if not fine_segments:
        return [text[:max_tokens * 4]] if text.strip() else []

    chunks: list[str] = []
    cur: list[str] = []
    cur_tok = 0
    for seg in fine_segments:
        st = count_tokens(seg)
        # If a single segment exceeds max_tokens, force-split by character
        if st > max_tokens:
            # Flush current chunk first
            if cur:
                chunks.append("".join(cur).strip())
                cur, cur_tok = [], 0
            # Split oversized segment by approximate char slices
            char_per_token = len(seg) / st if st > 0 else 4
            chunk_chars = int(max_tokens * char_per_token)
            for offset in range(0, len(seg), chunk_chars - int(overlap_tokens * char_per_token)):
                sub = seg[offset:offset + chunk_chars].strip()
                if sub:
                    chunks.append(sub)
            continue

        if cur_tok + st > max_tokens and cur:
            chunks.append("".join(cur).strip())
            ov: list[str] = []
            ov_t = 0
            for s in reversed(cur):
                t = count_tokens(s)
                if ov_t + t > overlap_tokens:
                    break
                ov.insert(0, s)
                ov_t += t
            cur = ov
            cur_tok = ov_t
        cur.append(seg)
        cur_tok += st
    if cur:
        f = "".join(cur).strip()
        if f:
            chunks.append(f)
    return chunks


# ---------- Orchestration ----------

async def process_file(file_id: UUID, db: AsyncSession) -> None:
    """Extract → chunk → embed → save (chunks + DocumentImage rows)."""
    from quip.routers.files import UPLOAD_DIR

    result = await db.execute(select(File).where(File.id == file_id))
    file_record = result.scalar_one_or_none()
    if not file_record:
        return

    file_record.embedding_status = "processing"
    await db.commit()

    try:
        file_path = UPLOAD_DIR / file_record.storage_path
        ext_result = await extract(str(file_path), file_record.content_type or "")

        # Archive: just one marker chunk, no embedding needed
        if ext_result.sandbox_only:
            sandbox_path = (
                f"/workspace/{file_record.chat_id}/uploads/{file_record.filename}"
                if file_record.chat_id
                else f"/workspace/<chat>/uploads/{file_record.filename}"
            )
            chunk = DocumentChunk(
                file_id=file_id,
                chat_id=file_record.chat_id,
                chunk_index=0,
                content=f"[archive: {file_record.filename} — available in sandbox at {sandbox_path}]",
                embedding=None,
                token_count=0,
                chunk_metadata={"source": "archive"},
            )
            db.add(chunk)
            file_record.embedding_status = "completed"
            await db.commit()
            logger.info(f"Archive {file_id}: marked sandbox_only")
            return

        if not ext_result.pages or not any(p.text.strip() for p in ext_result.pages):
            file_record.embedding_status = "failed"
            await db.commit()
            return

        # Persist extracted images on disk + DB
        if ext_result.images:
            await _save_images(file_record, ext_result.images, db)

        # Chunk
        max_tokens = int(get_setting("rag_chunk_size", "512"))
        overlap = int(get_setting("rag_chunk_overlap", "64"))
        chunks_with_meta = await asyncio.to_thread(
            chunk_pages, ext_result.pages, max_tokens=max_tokens, overlap_tokens=overlap
        )

        if not chunks_with_meta:
            file_record.embedding_status = "failed"
            await db.commit()
            return

        # Embed
        from quip.services.embeddings import get_embeddings
        texts = [c[0] for c in chunks_with_meta]
        embeddings = await get_embeddings(texts)

        if not embeddings or len(embeddings) != len(texts):
            file_record.embedding_status = "failed"
            await db.commit()
            return

        try:
            import tiktoken
            enc = tiktoken.get_encoding("cl100k_base")
            count_tokens = lambda t: len(enc.encode(t))
        except ImportError:
            count_tokens = lambda t: len(t) // 4

        for i, ((text, meta, chash), embedding) in enumerate(zip(chunks_with_meta, embeddings)):
            db.add(DocumentChunk(
                file_id=file_id,
                chat_id=file_record.chat_id,
                chunk_index=i,
                content=text,
                embedding=embedding,
                token_count=count_tokens(text),
                chunk_metadata=meta,
                content_hash=chash,
            ))

        file_record.embedding_status = "completed"
        await db.commit()
        logger.info(f"Processed file {file_id}: {len(chunks_with_meta)} chunks, {len(ext_result.images)} images")

    except Exception as e:
        logger.error(f"File processing failed for {file_id}: {e}")
        file_record.embedding_status = "failed"
        try:
            await db.commit()
        except Exception:
            pass


async def _save_images(file_record: File, images: list[ExtractedImage], db: AsyncSession) -> None:
    """Persist extracted images to disk + insert DocumentImage rows."""
    from quip.routers.files import UPLOAD_DIR

    user_dir = UPLOAD_DIR / str(file_record.user_id) / "extracted" / str(file_record.id)
    user_dir.mkdir(parents=True, exist_ok=True)

    ext_map = {
        "image/png": ".png",
        "image/jpeg": ".jpg",
        "image/jpg": ".jpg",
        "image/gif": ".gif",
        "image/webp": ".webp",
        "image/bmp": ".bmp",
        "image/tiff": ".tif",
    }

    for img in images:
        ext = ext_map.get(img.mime, ".bin")
        fname = f"{img.ref}{ext}"
        fpath = user_dir / fname
        try:
            fpath.write_bytes(img.data)
        except Exception as e:
            logger.warning(f"failed to write extracted image {img.ref}: {e}")
            continue
        storage_path = f"{file_record.user_id}/extracted/{file_record.id}/{fname}"
        db.add(DocumentImage(
            id=_uuid.uuid4(),
            file_id=file_record.id,
            ref=img.ref,
            page=img.page,
            storage_path=storage_path,
            mime=img.mime,
        ))


# ---------- Backwards compat shim ----------

async def extract_text(file_path: str, content_type: str) -> str:
    """Async text-only entrypoint. Returns concatenated page text."""
    res = await extract(file_path, content_type)
    return "\n\n".join(p.text for p in res.pages if p.text.strip())


def chunk_text(text: str, max_tokens: int = 512, overlap_tokens: int = 64) -> list[str]:
    """Legacy chunker — single-page text in, list of strings out."""
    chunks = chunk_pages([PageContent(text=text)], max_tokens=max_tokens, overlap_tokens=overlap_tokens)
    return [c[0] for c in chunks]
